from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from zipfile import BadZipFile, ZipFile

from docx import Document
from openpyxl import load_workbook
from pptx.enum.text import MSO_AUTO_SIZE
from pptx import Presentation
from pptx.util import Pt as PptxPt

from .models import ReplacementRule
from .text_replacer import plan_replacements, replace_text

LAYOUT_WARNING_PREFIX = "layout warning:"
_EMU_PER_POINT = 12700
_DEFAULT_PPTX_FONT_SIZE_PT = 12
_MIN_READABLE_FONT_SIZE_PT = 7
_LINE_HEIGHT_FACTOR = 1.2
_ASCII_CHARACTER_WIDTH_FACTOR = 0.55
_CJK_CHARACTER_WIDTH_FACTOR = 1.0


@dataclass(frozen=True)
class _TextSegment:
    text: str
    source_run: object


@dataclass(frozen=True)
class _PptxTextRegion:
    text_frame: object
    width: int
    height: int
    reference: str


def replace_docx(source: str | Path, output: str | Path, rules: tuple[ReplacementRule, ...]) -> tuple[int, list[str]]:
    document = Document(str(source))
    count = 0
    for paragraph in document.paragraphs:
        count += _replace_paragraph_runs(paragraph, rules)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    count += _replace_paragraph_runs(paragraph, rules)

    output_path = Path(output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(output_path))
    return count, embedded_object_notes(source)


def replace_xlsx(source: str | Path, output: str | Path, rules: tuple[ReplacementRule, ...]) -> tuple[int, list[str]]:
    workbook = load_workbook(source)
    count = 0
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    replaced, replacements = replace_text(cell.value, rules)
                    if replacements:
                        cell.value = replaced
                        count += replacements
    output_path = Path(output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output_path)
    return count, embedded_object_notes(source)


def replace_pptx(source: str | Path, output: str | Path, rules: tuple[ReplacementRule, ...]) -> tuple[int, list[str]]:
    presentation = Presentation(str(source))
    count = 0
    messages: list[str] = []
    source_name = Path(source).name
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                region = _shape_text_region(shape)
                replacements = _replace_pptx_region(region, rules)
                count += replacements
                if replacements and not _set_readable_text_frame_layout(region.text_frame, region.width, region.height):
                    messages.append(_layout_warning(source_name, slide_index, region.reference))
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for column_index, cell in enumerate(row.cells):
                        region = _table_cell_text_region(shape, row_index, column_index, cell)
                        replacements = _replace_pptx_region(region, rules)
                        count += replacements
                        if replacements and not _set_readable_text_frame_layout(region.text_frame, region.width, region.height):
                            messages.append(_layout_warning(source_name, slide_index, region.reference))
    output_path = Path(output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return count, messages + embedded_object_notes(source)


def embedded_object_notes(source: str | Path) -> list[str]:
    try:
        with ZipFile(source) as archive:
            names = archive.namelist()
    except (BadZipFile, FileNotFoundError):
        return []
    indicators = ("/embeddings/", "embeddings/", "/oleObject", "oleObject")
    if any(any(indicator in name for indicator in indicators) for name in names):
        return ["embedded objects are out of scope and were not masked"]
    return []


def _replace_paragraph_runs(paragraph, rules: tuple[ReplacementRule, ...]) -> int:
    return _replace_visible_run_text(paragraph, rules)


def _replace_pptx_runs(paragraph, rules: tuple[ReplacementRule, ...]) -> int:
    return _replace_visible_run_text(paragraph, rules)


def _replace_pptx_region(region: _PptxTextRegion, rules: tuple[ReplacementRule, ...]) -> int:
    count = 0
    for paragraph in region.text_frame.paragraphs:
        count += _replace_pptx_runs(paragraph, rules)
    return count


def _shape_text_region(shape) -> _PptxTextRegion:
    name = getattr(shape, "name", "") or "unnamed"
    return _PptxTextRegion(
        text_frame=shape.text_frame,
        width=int(getattr(shape, "width", 0) or 0),
        height=int(getattr(shape, "height", 0) or 0),
        reference=f"shape '{name}'",
    )


def _table_cell_text_region(shape, row_index: int, column_index: int, cell) -> _PptxTextRegion:
    table = shape.table
    try:
        width = int(table.columns[column_index].width)
    except (AttributeError, IndexError, TypeError):
        width = int(getattr(shape, "width", 0) or 0)
    try:
        height = int(table.rows[row_index].height)
    except (AttributeError, IndexError, TypeError):
        height = int(getattr(shape, "height", 0) or 0)
    name = getattr(shape, "name", "") or "table"
    return _PptxTextRegion(
        text_frame=cell.text_frame,
        width=width,
        height=height,
        reference=f"table '{name}' cell R{row_index + 1}C{column_index + 1}",
    )


def _layout_warning(source_name: str, slide_index: int, region_reference: str) -> str:
    return (
        f"{LAYOUT_WARNING_PREFIX} {source_name} slide {slide_index} {region_reference} "
        "may need manual review because replacement text cannot be kept readable inside the original region"
    )


def _set_readable_text_frame_layout(text_frame, width: int, height: int) -> bool:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    size_pt = _largest_readable_font_size(text_frame, width, height)
    if size_pt is None:
        _apply_pptx_font_size(text_frame, _MIN_READABLE_FONT_SIZE_PT)
        return False
    _apply_pptx_font_size(text_frame, size_pt)
    return True


def pptx_text_region_fits(text_frame, width: int, height: int) -> bool:
    size_pt = _smallest_declared_font_size(text_frame) or _DEFAULT_PPTX_FONT_SIZE_PT
    return _text_frame_fits_at_size(text_frame, width, height, size_pt)


def _largest_readable_font_size(text_frame, width: int, height: int) -> int | None:
    current_size = _largest_declared_font_size(text_frame) or _DEFAULT_PPTX_FONT_SIZE_PT
    for size_pt in range(current_size, _MIN_READABLE_FONT_SIZE_PT - 1, -1):
        if _text_frame_fits_at_size(text_frame, width, height, size_pt):
            return size_pt
    return None


def _text_frame_fits_at_size(text_frame, width: int, height: int, size_pt: float) -> bool:
    available_width = _available_text_width_pt(text_frame, width)
    available_height = _available_text_height_pt(text_frame, height)
    if available_width <= 0 or available_height <= 0:
        return False

    line_capacity = max(1.0, available_width / max(1.0, size_pt * _ASCII_CHARACTER_WIDTH_FACTOR))
    estimated_lines = 0
    for line in _text_frame_lines(text_frame):
        weighted_length = sum(_character_width_units(character) for character in line)
        estimated_lines += max(1, int((weighted_length + line_capacity - 1) // line_capacity))

    required_height = estimated_lines * size_pt * _LINE_HEIGHT_FACTOR
    return required_height <= available_height


def _available_text_width_pt(text_frame, width: int) -> float:
    margins = _length_to_pt(getattr(text_frame, "margin_left", 0)) + _length_to_pt(getattr(text_frame, "margin_right", 0))
    return max(0.0, _length_to_pt(width) - margins)


def _available_text_height_pt(text_frame, height: int) -> float:
    margins = _length_to_pt(getattr(text_frame, "margin_top", 0)) + _length_to_pt(getattr(text_frame, "margin_bottom", 0))
    return max(0.0, _length_to_pt(height) - margins)


def _length_to_pt(value) -> float:
    try:
        return float(int(value)) / _EMU_PER_POINT
    except (TypeError, ValueError):
        return 0.0


def _text_frame_lines(text_frame) -> list[str]:
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs)
        lines.extend(text.splitlines() or [text])
    return [line for line in lines if line] or [""]


def _character_width_units(character: str) -> float:
    return _ASCII_CHARACTER_WIDTH_FACTOR if ord(character) < 128 else _CJK_CHARACTER_WIDTH_FACTOR


def _apply_pptx_font_size(text_frame, size_pt: int) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = PptxPt(size_pt)


def _largest_declared_font_size(text_frame) -> int | None:
    sizes = _declared_font_sizes(text_frame)
    return max(sizes) if sizes else None


def _smallest_declared_font_size(text_frame) -> int | None:
    sizes = _declared_font_sizes(text_frame)
    return min(sizes) if sizes else None


def _declared_font_sizes(text_frame) -> list[int]:
    sizes: list[int] = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(round(_length_to_pt(run.font.size)))
    return sizes


def _replace_visible_run_text(paragraph, rules: tuple[ReplacementRule, ...]) -> int:
    runs = list(paragraph.runs)
    segments = [_TextSegment(run.text, run) for run in runs if run.text]
    if not segments:
        return 0

    visible_text = "".join(segment.text for segment in segments)
    matches = plan_replacements(visible_text, rules)
    if not matches:
        return 0

    next_segments: list[_TextSegment] = []
    cursor = 0
    for match in matches:
        next_segments.extend(_slice_segments(segments, cursor, match.start))
        first_run = _run_at(segments, match.start)
        next_segments.append(_TextSegment(match.replacement, first_run))
        cursor = match.end
    next_segments.extend(_slice_segments(segments, cursor, len(visible_text)))
    _write_segments(paragraph, next_segments)
    return len(matches)


def _slice_segments(segments: list[_TextSegment], start: int, end: int) -> list[_TextSegment]:
    if start >= end:
        return []

    sliced: list[_TextSegment] = []
    cursor = 0
    for segment in segments:
        segment_end = cursor + len(segment.text)
        overlap_start = max(start, cursor)
        overlap_end = min(end, segment_end)
        if overlap_start < overlap_end:
            local_start = overlap_start - cursor
            local_end = overlap_end - cursor
            sliced.append(_TextSegment(segment.text[local_start:local_end], segment.source_run))
        cursor = segment_end
    return sliced


def _run_at(segments: list[_TextSegment], index: int):
    cursor = 0
    for segment in segments:
        segment_end = cursor + len(segment.text)
        if cursor <= index < segment_end:
            return segment.source_run
        cursor = segment_end
    return segments[-1].source_run


def _write_segments(paragraph, segments: list[_TextSegment]) -> None:
    runs = list(paragraph.runs)
    while len(runs) < len(segments):
        runs.append(paragraph.add_run())

    for run, segment in zip(runs, segments):
        _copy_run_format(run, segment.source_run)
        run.text = segment.text

    for run in runs[len(segments) :]:
        run.text = ""


def _copy_run_format(target, source) -> None:
    if hasattr(target, "bold"):
        target.bold = getattr(source, "bold", None)
    if hasattr(target, "italic"):
        target.italic = getattr(source, "italic", None)

    if hasattr(target, "font") and hasattr(source, "font"):
        _copy_font_format(target.font, source.font)


def _copy_font_format(target_font, source_font) -> None:
    for attr in ("bold", "italic", "size"):
        if hasattr(target_font, attr) and hasattr(source_font, attr):
            setattr(target_font, attr, getattr(source_font, attr))

    try:
        source_rgb = source_font.color.rgb
    except AttributeError:
        source_rgb = None
    if source_rgb is not None:
        try:
            target_font.color.rgb = source_rgb
        except AttributeError:
            pass
