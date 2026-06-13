from __future__ import annotations

from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

from masking_tool.models import ReplacementRule
from masking_tool.office_replacer import (
    LAYOUT_WARNING_PREFIX,
    _layout_warning,
    _replace_paragraph_runs,
    _set_readable_text_frame_layout,
    pptx_text_region_fits,
)
from tests.fixtures.create_fixtures import SPLIT_PHRASE, SPLIT_REPL
from tests.fixtures.create_replacement_layout_samples import LAYOUT_LONG_REPL


def test_paragraph_helper_replaces_visible_split_phrase_with_first_run_format():
    document = Document()
    paragraph = document.add_paragraph()
    first = paragraph.add_run("Technologies")
    first.bold = True
    first.italic = True
    first.font.size = Pt(16)
    first.font.color.rgb = RGBColor(0x12, 0x34, 0x56)
    paragraph.add_run(", Inc.")

    count = _replace_paragraph_runs(paragraph, (_rule(),))

    assert count == 1
    assert paragraph.text == SPLIT_REPL
    assert paragraph.runs[0].bold is True
    assert paragraph.runs[0].italic is True
    assert paragraph.runs[0].font.size.pt == 16
    assert str(paragraph.runs[0].font.color.rgb) == "123456"


def test_paragraph_helper_does_not_count_internal_run_fragments():
    document = Document()
    paragraph = document.add_paragraph()
    paragraph.add_run("Technologies")
    paragraph.add_run(", Inc.")
    paragraph.add_run(" / Technologies")
    paragraph.add_run(", Inc.")

    count = _replace_paragraph_runs(paragraph, (_rule(),))

    assert count == 2
    assert paragraph.text.count(SPLIT_REPL) == 2


def test_pptx_text_frame_layout_helper_reduces_font_to_fit_region():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, Inches(2.5), Inches(0.7))
    shape.text_frame.text = "担当者〇1 担当者〇2 担当者〇3"
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = PptxPt(18)

    readable = _set_readable_text_frame_layout(shape.text_frame, shape.width, shape.height)

    assert readable is True
    assert pptx_text_region_fits(shape.text_frame, shape.width, shape.height)


def test_pptx_text_frame_layout_helper_reports_unreadable_region():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, Inches(0.7), Inches(0.2))
    shape.text_frame.text = LAYOUT_LONG_REPL

    readable = _set_readable_text_frame_layout(shape.text_frame, shape.width, shape.height)

    assert readable is False


def test_layout_warning_includes_file_slide_and_region_context():
    message = _layout_warning("sample.pptx", 2, "shape 'owner'")

    assert message.startswith(LAYOUT_WARNING_PREFIX)
    assert "sample.pptx" in message
    assert "slide 2" in message
    assert "shape 'owner'" in message


def _rule() -> ReplacementRule:
    return ReplacementRule(no="1", detected_phrase=SPLIT_PHRASE, replacement_proposal=SPLIT_REPL, row_index=2)
