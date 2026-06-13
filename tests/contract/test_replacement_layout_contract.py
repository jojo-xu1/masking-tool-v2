from __future__ import annotations

from pptx import Presentation

from masking_tool.app import process_selection
from masking_tool.models import InputMode, InputSelection
from masking_tool.office_replacer import LAYOUT_WARNING_PREFIX, pptx_text_region_fits
from tests.fixtures.create_replacement_layout_samples import (
    LAYOUT_LONG_PHRASE,
    LAYOUT_LONG_REPL,
    LAYOUT_PHRASE_1,
    LAYOUT_PHRASE_2,
    LAYOUT_PHRASE_3,
    LAYOUT_PHRASE_4,
    LAYOUT_REPL_1,
    LAYOUT_REPL_2,
    LAYOUT_REPL_3,
    LAYOUT_REPL_4,
    create_layout_diagram_pptx,
    create_layout_overflow_pptx,
    create_layout_replacement_table,
    create_layout_textbox_and_table_pptx,
)


def test_screenshot_style_pptx_diagram_replacements_are_readable_without_original_phrases(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_diagram_pptx(tmp_path / "diagram.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 4
    assert not summary.results[0].messages
    presentation = Presentation(str(tmp_path / "out" / "diagram.pptx"))
    shape = presentation.slides[0].shapes[0]
    text = shape.text
    for phrase in (LAYOUT_PHRASE_1, LAYOUT_PHRASE_2, LAYOUT_PHRASE_3, LAYOUT_PHRASE_4):
        assert phrase not in text
    for replacement in (LAYOUT_REPL_1, LAYOUT_REPL_2, LAYOUT_REPL_3, LAYOUT_REPL_4):
        assert replacement in text
    assert pptx_text_region_fits(shape.text_frame, shape.width, shape.height)


def test_pptx_textbox_and_table_cell_replacements_are_readable(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_textbox_and_table_pptx(tmp_path / "textbox-table.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 4
    assert not summary.results[0].messages
    presentation = Presentation(str(tmp_path / "out" / "textbox-table.pptx"))
    textbox, table_shape = presentation.slides[0].shapes
    assert pptx_text_region_fits(textbox.text_frame, textbox.width, textbox.height)
    cell = table_shape.table.cell(0, 0)
    assert pptx_text_region_fits(cell.text_frame, table_shape.table.columns[0].width, table_shape.table.rows[0].height)


def test_unreadable_pptx_region_records_layout_warning(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_overflow_pptx(tmp_path / "overflow.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 1
    assert any(message.startswith(LAYOUT_WARNING_PREFIX) for message in summary.results[0].messages)
    presentation = Presentation(str(tmp_path / "out" / "overflow.pptx"))
    text = presentation.slides[0].shapes[0].text
    assert LAYOUT_LONG_PHRASE not in text
    assert LAYOUT_LONG_REPL in text
