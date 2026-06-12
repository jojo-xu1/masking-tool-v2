from __future__ import annotations

from docx import Document
from pptx import Presentation

from masking_tool.app import process_selection
from masking_tool.models import InputMode, InputSelection
from tests.fixtures.create_fixtures import (
    SPLIT_PHRASE,
    SPLIT_REPL,
    create_replacement_table,
    create_split_run_docx,
    create_split_run_pptx,
)


def test_docx_split_run_paragraph_replacement_uses_first_run_format_and_count(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, SPLIT_PHRASE, SPLIT_REPL)])
    source = create_split_run_docx(tmp_path / "sample.docx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 1
    doc = Document(str(tmp_path / "out" / "sample.docx"))
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    assert SPLIT_PHRASE not in text
    assert SPLIT_REPL in text
    replacement_run = _docx_run_containing(doc.paragraphs[0].runs, SPLIT_REPL)
    assert replacement_run.bold is True
    assert replacement_run.italic is True
    assert replacement_run.font.size.pt == 16
    assert str(replacement_run.font.color.rgb) == "123456"


def test_docx_split_run_table_cell_replacement(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, SPLIT_PHRASE, SPLIT_REPL)])
    source = create_split_run_docx(tmp_path / "sample.docx", in_table=True)

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 1
    doc = Document(str(tmp_path / "out" / "sample.docx"))
    cell_text = doc.tables[0].cell(0, 0).text
    assert SPLIT_PHRASE not in cell_text
    assert SPLIT_REPL in cell_text


def test_pptx_split_run_text_box_replacement_uses_first_run_format_and_count(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, SPLIT_PHRASE, SPLIT_REPL)])
    source = create_split_run_pptx(tmp_path / "sample.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 1
    prs = Presentation(str(tmp_path / "out" / "sample.pptx"))
    paragraph = prs.slides[0].shapes[1].text_frame.paragraphs[0]
    text = "".join(run.text for run in paragraph.runs)
    assert SPLIT_PHRASE not in text
    assert SPLIT_REPL in text
    replacement_run = _pptx_run_containing(paragraph.runs, SPLIT_REPL)
    assert replacement_run.font.bold is True
    assert replacement_run.font.italic is True
    assert replacement_run.font.size.pt == 16
    assert str(replacement_run.font.color.rgb) == "123456"


def test_pptx_split_run_table_cell_replacement(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, SPLIT_PHRASE, SPLIT_REPL)])
    source = create_split_run_pptx(tmp_path / "sample.pptx", in_table=True)

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 1
    prs = Presentation(str(tmp_path / "out" / "sample.pptx"))
    table_shape = prs.slides[0].shapes[1]
    cell_text = table_shape.table.cell(0, 0).text
    assert SPLIT_PHRASE not in cell_text
    assert SPLIT_REPL in cell_text


def test_no_match_split_run_adjacent_office_files_do_not_report_replacement(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, SPLIT_PHRASE, SPLIT_REPL)])
    docx = create_split_run_docx(tmp_path / "sample.docx", text="Technologies Inc")
    pptx = create_split_run_pptx(tmp_path / "sample.pptx", text="Technologies Inc")

    docx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, docx, tmp_path / "out-docx"))
    pptx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, pptx, tmp_path / "out-pptx"))

    assert docx_summary.results[0].replacement_count == 0
    assert pptx_summary.results[0].replacement_count == 0


def _docx_run_containing(runs, text):
    return next(run for run in runs if text in run.text)


def _pptx_run_containing(runs, text):
    return next(run for run in runs if text in run.text)
