from __future__ import annotations

from pptx import Presentation

from masking_tool.app import process_selection
from masking_tool.models import InputMode, InputSelection, TraversalMode
from masking_tool.office_replacer import LAYOUT_WARNING_PREFIX, pptx_text_region_fits
from tests.fixtures.create_replacement_layout_samples import (
    LAYOUT_LONG_REPL,
    LAYOUT_PHRASE_1,
    LAYOUT_PHRASE_2,
    LAYOUT_PHRASE_3,
    LAYOUT_PHRASE_4,
    LAYOUT_REPL_1,
    LAYOUT_REPL_2,
    LAYOUT_REPL_3,
    LAYOUT_REPL_4,
    create_layout_diagram_pptx,
    create_layout_mixed_pptx,
    create_layout_overflow_pptx,
    create_layout_replacement_table,
    create_layout_textbox_and_table_pptx,
)


def test_diagram_like_pptx_replacement_labels_do_not_overlap(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_diagram_pptx(tmp_path / "project.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.replaced_count == 1
    assert summary.results[0].replacement_count == 4
    assert not summary.results[0].messages
    shape = Presentation(str(tmp_path / "out" / "project.pptx")).slides[0].shapes[0]
    text = shape.text
    for phrase in (LAYOUT_PHRASE_1, LAYOUT_PHRASE_2, LAYOUT_PHRASE_3, LAYOUT_PHRASE_4):
        assert phrase not in text
    for replacement in (LAYOUT_REPL_1, LAYOUT_REPL_2, LAYOUT_REPL_3, LAYOUT_REPL_4):
        assert replacement in text
    assert pptx_text_region_fits(shape.text_frame, shape.width, shape.height)


def test_textbox_and_table_cell_layouts_fit_original_regions(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_textbox_and_table_pptx(tmp_path / "regions.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.results[0].replacement_count == 4
    assert not summary.results[0].messages
    presentation = Presentation(str(tmp_path / "out" / "regions.pptx"))
    textbox, table_shape = presentation.slides[0].shapes
    assert pptx_text_region_fits(textbox.text_frame, textbox.width, textbox.height)
    assert pptx_text_region_fits(
        table_shape.table.cell(0, 0).text_frame,
        table_shape.table.columns[0].width,
        table_shape.table.rows[0].height,
    )


def test_unreadable_layout_warns_in_result_and_skip_report(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_overflow_pptx(tmp_path / "overflow.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    report = summary.report_path.read_text(encoding="utf-8")
    assert summary.replaced_count == 1
    assert summary.results[0].replacement_count == 1
    assert any(message.startswith(LAYOUT_WARNING_PREFIX) for message in summary.results[0].messages)
    assert "overflow.pptx\tskipped_unsupported\tlayout warning:" in report
    assert LAYOUT_LONG_REPL in Presentation(str(tmp_path / "out" / "overflow.pptx")).slides[0].shapes[0].text


def test_mixed_readable_and_warning_regions_keep_readable_replacements_successful(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    source = create_layout_mixed_pptx(tmp_path / "mixed.pptx")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    assert summary.replaced_count == 1
    assert summary.results[0].replacement_count == 3
    assert any(message.startswith(LAYOUT_WARNING_PREFIX) for message in summary.results[0].messages)
    presentation = Presentation(str(tmp_path / "out" / "mixed.pptx"))
    readable = presentation.slides[0].shapes[0]
    warning = presentation.slides[0].shapes[1]
    assert LAYOUT_REPL_1 in readable.text
    assert LAYOUT_REPL_2 in readable.text
    assert pptx_text_region_fits(readable.text_frame, readable.width, readable.height)
    assert LAYOUT_LONG_REPL in warning.text


def test_folder_skip_report_records_layout_warning(tmp_path):
    table = create_layout_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    folder = tmp_path / "in"
    create_layout_overflow_pptx(folder / "overflow.pptx")

    summary = process_selection(
        table,
        InputSelection(InputMode.FOLDER, folder, tmp_path / "out", TraversalMode.DIRECT_CHILDREN),
    )

    report = summary.report_path.read_text(encoding="utf-8")
    assert summary.replaced_count == 1
    assert "overflow.pptx\tskipped_unsupported\tlayout warning:" in report
