from __future__ import annotations

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from masking_tool.app import process_selection
from masking_tool.models import InputMode, InputSelection
from tests.fixtures.create_fixtures import (
    UNICODE_FULL_WIDTH_PHRASE,
    UNICODE_FULL_WIDTH_REPL,
    UNICODE_HALF_WIDTH_PHRASE,
    UNICODE_REPL,
    create_docx,
    create_pptx,
    create_replacement_table,
    create_text_file,
    create_text_pdf,
    create_xlsx,
)


def test_text_formats_full_width_rule_masks_half_width_targets(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, UNICODE_FULL_WIDTH_PHRASE, UNICODE_REPL)])
    for suffix in [".txt", ".csv", ".log"]:
        source = create_text_file(tmp_path / f"input{suffix}", f"取引先: {UNICODE_HALF_WIDTH_PHRASE}")
        out = tmp_path / f"out{suffix[1:]}"

        summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, out))

        output_text = (out / source.name).read_text(encoding="utf-8")
        assert summary.results[0].replacement_count == 1
        assert UNICODE_HALF_WIDTH_PHRASE not in output_text
        assert UNICODE_REPL in output_text


def test_office_and_workbook_full_width_rule_masks_half_width_targets(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, UNICODE_FULL_WIDTH_PHRASE, UNICODE_REPL)])

    docx = create_docx(tmp_path / "sample.docx", f"取引先: {UNICODE_HALF_WIDTH_PHRASE}")
    docx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, docx, tmp_path / "out-docx"))
    assert docx_summary.results[0].replacement_count == 1
    assert UNICODE_REPL in _docx_text(tmp_path / "out-docx" / "sample.docx")
    assert UNICODE_HALF_WIDTH_PHRASE not in _docx_text(tmp_path / "out-docx" / "sample.docx")

    xlsx = create_xlsx(tmp_path / "sample.xlsx", f"取引先: {UNICODE_HALF_WIDTH_PHRASE}")
    xlsx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, xlsx, tmp_path / "out-xlsx"))
    workbook = load_workbook(tmp_path / "out-xlsx" / "sample.xlsx")
    assert xlsx_summary.results[0].replacement_count == 1
    assert workbook.active["A1"].value == f"取引先: {UNICODE_REPL}"

    pptx = create_pptx(tmp_path / "sample.pptx", f"取引先: {UNICODE_HALF_WIDTH_PHRASE}")
    pptx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, pptx, tmp_path / "out-pptx"))
    assert pptx_summary.results[0].replacement_count == 1
    assert UNICODE_REPL in _pptx_text(tmp_path / "out-pptx" / "sample.pptx")
    assert UNICODE_HALF_WIDTH_PHRASE not in _pptx_text(tmp_path / "out-pptx" / "sample.pptx")


def test_pdf_full_width_rule_masks_half_width_target(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, UNICODE_FULL_WIDTH_PHRASE, UNICODE_REPL)])
    source = create_text_pdf(tmp_path / "sample.pdf", f"取引先: {UNICODE_HALF_WIDTH_PHRASE}")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    text = _pdf_text(tmp_path / "out" / "sample.pdf")
    assert summary.results[0].replacement_count == 1
    assert UNICODE_HALF_WIDTH_PHRASE not in text
    assert UNICODE_REPL in text


def test_text_formats_half_width_rule_masks_full_width_targets(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, UNICODE_HALF_WIDTH_PHRASE, UNICODE_REPL)])
    for suffix in [".txt", ".csv", ".log"]:
        source = create_text_file(tmp_path / f"input{suffix}", f"取引先: {UNICODE_FULL_WIDTH_PHRASE}")
        out = tmp_path / f"out{suffix[1:]}"

        summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, out))

        output_text = (out / source.name).read_text(encoding="utf-8")
        assert summary.results[0].replacement_count == 1
        assert UNICODE_FULL_WIDTH_PHRASE not in output_text
        assert UNICODE_REPL in output_text


def test_office_and_workbook_half_width_rule_masks_full_width_targets(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, UNICODE_HALF_WIDTH_PHRASE, UNICODE_REPL)])

    docx = create_docx(tmp_path / "sample.docx", f"取引先: {UNICODE_FULL_WIDTH_PHRASE}")
    process_selection(table, InputSelection(InputMode.SINGLE_FILE, docx, tmp_path / "out-docx"))
    assert UNICODE_REPL in _docx_text(tmp_path / "out-docx" / "sample.docx")
    assert UNICODE_FULL_WIDTH_PHRASE not in _docx_text(tmp_path / "out-docx" / "sample.docx")

    xlsx = create_xlsx(tmp_path / "sample.xlsx", f"取引先: {UNICODE_FULL_WIDTH_PHRASE}")
    process_selection(table, InputSelection(InputMode.SINGLE_FILE, xlsx, tmp_path / "out-xlsx"))
    workbook = load_workbook(tmp_path / "out-xlsx" / "sample.xlsx")
    assert workbook.active["A1"].value == f"取引先: {UNICODE_REPL}"

    pptx = create_pptx(tmp_path / "sample.pptx", f"取引先: {UNICODE_FULL_WIDTH_PHRASE}")
    process_selection(table, InputSelection(InputMode.SINGLE_FILE, pptx, tmp_path / "out-pptx"))
    assert UNICODE_REPL in _pptx_text(tmp_path / "out-pptx" / "sample.pptx")
    assert UNICODE_FULL_WIDTH_PHRASE not in _pptx_text(tmp_path / "out-pptx" / "sample.pptx")


def test_pdf_half_width_rule_masks_full_width_target(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, UNICODE_HALF_WIDTH_PHRASE, UNICODE_REPL)])
    source = create_text_pdf(tmp_path / "sample.pdf", f"取引先: {UNICODE_FULL_WIDTH_PHRASE}")

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    text = _pdf_text(tmp_path / "out" / "sample.pdf")
    assert summary.results[0].replacement_count == 1
    assert UNICODE_FULL_WIDTH_PHRASE not in text
    assert UNICODE_REPL in text


def test_exact_raw_match_precedence_and_replacement_proposal_preservation(tmp_path):
    table = create_replacement_table(
        tmp_path / "機密情報検出結果.xlsx",
        [
            (1, UNICODE_FULL_WIDTH_PHRASE, "WIDTH_EQUIVALENT"),
            (2, UNICODE_HALF_WIDTH_PHRASE, UNICODE_FULL_WIDTH_REPL),
        ],
    )
    source = create_text_file(tmp_path / "input.txt", f"東京本社 {UNICODE_HALF_WIDTH_PHRASE} 完了")

    process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    output_text = (tmp_path / "out" / "input.txt").read_text(encoding="utf-8")
    assert output_text == f"東京本社 {UNICODE_FULL_WIDTH_REPL} 完了"


def _docx_text(path):
    document = Document(str(path))
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.extend(paragraph.text for paragraph in cell.paragraphs)
    return "\n".join(parts)


def _pptx_text(path):
    presentation = Presentation(str(path))
    parts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


def _pdf_text(path):
    document = fitz.open(path)
    try:
        return "\n".join(page.get_text("text") for page in document)
    finally:
        document.close()
