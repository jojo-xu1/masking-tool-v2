from __future__ import annotations

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt

from masking_tool.app import process_selection
from masking_tool.models import InputMode, InputSelection
from tests.fixtures.create_fixtures import (
    PHRASE_1,
    REPL_1,
    create_docx,
    create_pptx,
    create_replacement_table,
    create_xlsx,
)


def test_single_docx_xlsx_pptx_files(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx")

    docx = create_docx(tmp_path / "sample.docx")
    docx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, docx, tmp_path / "out-docx"))
    doc = Document(str(tmp_path / "out-docx" / "sample.docx"))
    assert docx_summary.results[0].replacement_count == 2
    assert PHRASE_1 not in "\n".join(p.text for p in doc.paragraphs)
    assert REPL_1 in "\n".join(p.text for p in doc.paragraphs)

    xlsx = create_xlsx(tmp_path / "sample.xlsx")
    process_selection(table, InputSelection(InputMode.SINGLE_FILE, xlsx, tmp_path / "out-xlsx"))
    workbook = load_workbook(tmp_path / "out-xlsx" / "sample.xlsx")
    assert workbook.active["A1"].value == REPL_1

    pptx = create_pptx(tmp_path / "sample.pptx")
    pptx_summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, pptx, tmp_path / "out-pptx"))
    prs = Presentation(str(tmp_path / "out-pptx" / "sample.pptx"))
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert pptx_summary.results[0].replacement_count == 2
    assert PHRASE_1 not in text
    assert REPL_1 in text


def test_pptx_replacement_preserves_mixed_font_sizes_in_same_textbox(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, "会社04", "会社XX")])
    source = tmp_path / "mixed-fonts.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, 8000000, 1600000)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.clear()
    for text, size in (("From ", 18), ("会社04", 16), (" to Medical Industry Standards.", 14)):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
    presentation.save(source)

    summary = process_selection(table, InputSelection(InputMode.SINGLE_FILE, source, tmp_path / "out"))

    output = Presentation(str(tmp_path / "out" / "mixed-fonts.pptx"))
    runs = output.slides[0].shapes[0].text_frame.paragraphs[0].runs
    assert summary.results[0].replacement_count == 1
    assert [run.text for run in runs] == ["From ", "会社XX", " to Medical Industry Standards."]
    assert [round(run.font.size.pt) for run in runs] == [18, 16, 14]
