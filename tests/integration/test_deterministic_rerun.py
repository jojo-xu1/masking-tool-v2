from __future__ import annotations

from docx import Document
from pptx import Presentation

from masking_tool.app import process_selection
from masking_tool.models import InputMode, InputSelection, TraversalMode
from tests.fixtures.create_fixtures import (
    SPLIT_PHRASE,
    SPLIT_REPL,
    create_replacement_table,
    create_split_run_docx,
    create_split_run_pptx,
    create_text_file,
)


def test_deterministic_rerun_outputs_and_reports(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx")
    folder = tmp_path / "in"
    create_text_file(folder / "a.txt")
    (folder / "skip.bin").write_text("x", encoding="utf-8")

    first = process_selection(table, InputSelection(InputMode.FOLDER, folder, tmp_path / "out1", TraversalMode.DIRECT_CHILDREN))
    second = process_selection(table, InputSelection(InputMode.FOLDER, folder, tmp_path / "out2", TraversalMode.DIRECT_CHILDREN))

    assert (tmp_path / "out1" / "a.txt").read_text(encoding="utf-8") == (tmp_path / "out2" / "a.txt").read_text(encoding="utf-8")
    assert first.report_path.read_text(encoding="utf-8") == second.report_path.read_text(encoding="utf-8")


def test_deterministic_rerun_split_run_office_outputs_and_counts(tmp_path):
    table = create_replacement_table(tmp_path / "機密情報検出結果.xlsx", [(1, SPLIT_PHRASE, SPLIT_REPL)])
    docx = create_split_run_docx(tmp_path / "sample.docx")
    pptx = create_split_run_pptx(tmp_path / "sample.pptx")

    first_docx = process_selection(table, InputSelection(InputMode.SINGLE_FILE, docx, tmp_path / "out1-docx"))
    second_docx = process_selection(table, InputSelection(InputMode.SINGLE_FILE, docx, tmp_path / "out2-docx"))
    first_pptx = process_selection(table, InputSelection(InputMode.SINGLE_FILE, pptx, tmp_path / "out1-pptx"))
    second_pptx = process_selection(table, InputSelection(InputMode.SINGLE_FILE, pptx, tmp_path / "out2-pptx"))

    assert first_docx.results[0].replacement_count == second_docx.results[0].replacement_count == 1
    assert first_pptx.results[0].replacement_count == second_pptx.results[0].replacement_count == 1
    assert _docx_text(tmp_path / "out1-docx" / "sample.docx") == _docx_text(tmp_path / "out2-docx" / "sample.docx")
    assert _pptx_text(tmp_path / "out1-pptx" / "sample.pptx") == _pptx_text(tmp_path / "out2-pptx" / "sample.pptx")


def _docx_text(path):
    doc = Document(str(path))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def _pptx_text(path):
    prs = Presentation(str(path))
    return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
