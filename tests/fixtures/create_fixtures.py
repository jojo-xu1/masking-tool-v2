from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from docx.shared import Pt, RGBColor
from openpyxl import Workbook
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Pt as PptxPt

PHRASE_1 = "山田太郎"
PHRASE_2 = "example@example.com"
SPLIT_PHRASE = "Technologies, Inc."
REPL_1 = "<PERSON_001>"
REPL_2 = "<EMAIL_001>"
SPLIT_REPL = "会社名_置換済み"


def create_replacement_table(path: Path, rows: list[tuple[object, str, str]] | None = None) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["No", "検出語句", "置換提案"])
    for row in rows or [(1, PHRASE_1, REPL_1), (2, PHRASE_2, REPL_2)]:
        sheet.append(list(row))
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(path)
    return path


def create_invalid_replacement_table(path: Path, headers: list[str]) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(headers)
    sheet.append([1, PHRASE_1])
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(path)
    return path


def create_text_file(path: Path, text: str | None = None) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text if text is not None else f"{PHRASE_1} / {PHRASE_2}", encoding="utf-8")
    return path


def create_docx(path: Path, text: str | None = None) -> Path:
    document = Document()
    document.add_paragraph(text if text is not None else f"担当者: {PHRASE_1}")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = PHRASE_2
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)
    return path


def create_split_run_docx(path: Path, *, in_table: bool = False, text: str | None = None) -> Path:
    document = Document()
    paragraph = document.add_paragraph("概要: ")
    if in_table:
        table = document.add_table(rows=1, cols=1)
        paragraph = table.cell(0, 0).paragraphs[0]
        paragraph.text = ""
    _add_docx_split_runs(paragraph, text or SPLIT_PHRASE)
    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)
    return path


def create_xlsx(path: Path, text: str | None = None) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = text if text is not None else PHRASE_1
    sheet["B1"] = PHRASE_2
    path.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(path)
    return path


def create_pptx(path: Path, text: str | None = None) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 4000000, 1000000)
    box.text_frame.text = text if text is not None else f"{PHRASE_1} {PHRASE_2}"
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def create_split_run_pptx(path: Path, *, in_table: bool = False, text: str | None = None) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    if in_table:
        shape = slide.shapes.add_table(1, 1, 0, 0, 4000000, 1000000)
        paragraph = shape.table.cell(0, 0).text_frame.paragraphs[0]
    else:
        shape = slide.shapes.add_textbox(0, 0, 4000000, 1000000)
        paragraph = shape.text_frame.paragraphs[0]
    _add_pptx_split_runs(paragraph, text or SPLIT_PHRASE)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def create_text_pdf(path: Path, text: str | None = None) -> Path:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text if text is not None else PHRASE_2)
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(path)
    doc.close()
    return path


def create_scanned_like_pdf(path: Path) -> Path:
    doc = fitz.open()
    doc.new_page()
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(path)
    doc.close()
    return path


def _add_docx_split_runs(paragraph, text: str) -> None:
    paragraph.add_run("取引先: ")
    if text == SPLIT_PHRASE:
        first = "Technologies"
        second = ", Inc."
    else:
        first = text[: max(1, len(text) // 2)]
        second = text[len(first) :]
    first_run = paragraph.add_run(first)
    first_run.bold = True
    first_run.italic = True
    first_run.font.size = Pt(16)
    first_run.font.color.rgb = RGBColor(0x12, 0x34, 0x56)
    second_run = paragraph.add_run(second)
    second_run.bold = False
    second_run.italic = False
    second_run.font.size = Pt(9)
    second_run.font.color.rgb = RGBColor(0xAA, 0x22, 0x22)
    paragraph.add_run(" / 完了")


def _add_pptx_split_runs(paragraph, text: str) -> None:
    prefix = paragraph.add_run()
    prefix.text = "取引先: "
    if text == SPLIT_PHRASE:
        first = "Technologies"
        second = ", Inc."
    else:
        first = text[: max(1, len(text) // 2)]
        second = text[len(first) :]
    first_run = paragraph.add_run()
    first_run.text = first
    first_run.font.bold = True
    first_run.font.italic = True
    first_run.font.size = PptxPt(16)
    first_run.font.color.rgb = PptxRGBColor(0x12, 0x34, 0x56)
    second_run = paragraph.add_run()
    second_run.text = second
    second_run.font.bold = False
    second_run.font.italic = False
    second_run.font.size = PptxPt(9)
    second_run.font.color.rgb = PptxRGBColor(0xAA, 0x22, 0x22)
    suffix = paragraph.add_run()
    suffix.text = " / 完了"
