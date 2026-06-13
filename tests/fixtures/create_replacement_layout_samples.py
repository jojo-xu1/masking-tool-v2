from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .create_fixtures import create_replacement_table

LAYOUT_PHRASE_1 = "担当者A"
LAYOUT_PHRASE_2 = "担当者B"
LAYOUT_PHRASE_3 = "担当者C"
LAYOUT_PHRASE_4 = "会社O2"
LAYOUT_REPL_1 = "担当者〇1"
LAYOUT_REPL_2 = "担当者〇2"
LAYOUT_REPL_3 = "担当者〇3"
LAYOUT_REPL_4 = "会社〇2"
LAYOUT_LONG_PHRASE = "承認者"
LAYOUT_LONG_REPL = "非常に長い置換ラベル_手動確認が必要_001"

LAYOUT_ROWS = [
    (1, LAYOUT_PHRASE_1, LAYOUT_REPL_1),
    (2, LAYOUT_PHRASE_2, LAYOUT_REPL_2),
    (3, LAYOUT_PHRASE_3, LAYOUT_REPL_3),
    (4, LAYOUT_PHRASE_4, LAYOUT_REPL_4),
    (5, LAYOUT_LONG_PHRASE, LAYOUT_LONG_REPL),
]
ROOT = Path(__file__).resolve().parent
INPUT_DIR = ROOT / "inputs" / "replacement_layout_samples"
TABLE_DIR = ROOT / "replacement_tables"


def create_layout_replacement_table(path: Path) -> Path:
    return create_replacement_table(path, LAYOUT_ROWS)


def create_layout_diagram_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(0.6), Inches(4.4), Inches(1.0))
    shape.name = "project-summary"
    shape.text_frame.clear()
    shape.text_frame.margin_left = Inches(0.08)
    shape.text_frame.margin_right = Inches(0.08)
    shape.text_frame.margin_top = Inches(0.05)
    shape.text_frame.margin_bottom = Inches(0.05)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = f"{LAYOUT_PHRASE_1} {LAYOUT_PHRASE_2} {LAYOUT_PHRASE_3} {LAYOUT_PHRASE_4}"
    for run in paragraph.runs:
        run.font.size = Pt(16)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def create_layout_textbox_and_table_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(3.6), Inches(0.8))
    box.name = "owner-textbox"
    box.text_frame.text = f"{LAYOUT_PHRASE_1}: {LAYOUT_PHRASE_2}"
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    table_shape = slide.shapes.add_table(1, 1, Inches(0.6), Inches(1.6), Inches(3.4), Inches(0.8))
    table_shape.name = "owner-table"
    cell = table_shape.table.cell(0, 0)
    cell.text = f"{LAYOUT_PHRASE_3} / {LAYOUT_PHRASE_4}"
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def create_layout_overflow_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(0.8), Inches(0.25))
    box.name = "too-small-textbox"
    box.text_frame.text = LAYOUT_LONG_PHRASE
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def create_layout_mixed_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    readable = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3.6), Inches(0.8))
    readable.name = "readable-textbox"
    readable.text_frame.text = f"{LAYOUT_PHRASE_1} {LAYOUT_PHRASE_2}"

    warning = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(0.8), Inches(0.25))
    warning.name = "warning-textbox"
    warning.text_frame.text = LAYOUT_LONG_PHRASE

    for shape in (readable, warning):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def create_replacement_layout_inputs() -> None:
    INPUT_DIR.mkdir(parents=True, exist_ok=True)
    create_layout_diagram_pptx(INPUT_DIR / "layout_diagram.pptx")
    create_layout_textbox_and_table_pptx(INPUT_DIR / "layout_textbox_table.pptx")
    create_layout_overflow_pptx(INPUT_DIR / "layout_overflow.pptx")
    create_layout_mixed_pptx(INPUT_DIR / "layout_mixed.pptx")


def create_replacement_layout_table() -> Path:
    return create_layout_replacement_table(TABLE_DIR / "機密情報検出結果_replacement_layout.xlsx")


def main() -> None:
    create_replacement_layout_table()
    create_replacement_layout_inputs()


if __name__ == "__main__":
    main()
