from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

ROOT = Path(__file__).resolve().parent
INPUT_DIR = ROOT / "inputs" / "manual_sensitive_samples"
TABLE_DIR = ROOT / "replacement_tables"

PHRASES = [
    "増田",
    "ＴＩ推",
    "HAND趙",
    "0125-74-2111",
]

REPLACEMENTS = [
    "<PERSON_001>",
    "<ORG_001>",
    "<PERSON_002>",
    "<PHONE_001>",
]


def sample_text() -> str:
    return "\n".join(
        [
            "機密情報マスキングテスト",
            f"担当者: {PHRASES[0]}",
            f"所属: {PHRASES[1]}",
            f"担当コード: {PHRASES[2]}",
            f"電話番号: {PHRASES[3]}",
        ]
    )


def create_replacement_table() -> Path:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "検出結果"
    sheet.append(["No", "検出語句", "置換提案"])
    for index, (phrase, replacement) in enumerate(zip(PHRASES, REPLACEMENTS), start=1):
        sheet.append([index, phrase, replacement])
    TABLE_DIR.mkdir(parents=True, exist_ok=True)
    path = TABLE_DIR / "機密情報検出結果_manual_sensitive.xlsx"
    workbook.save(path)
    return path


def create_text_files() -> None:
    INPUT_DIR.mkdir(parents=True, exist_ok=True)
    text = sample_text()
    (INPUT_DIR / "manual_sample.txt").write_text(text + "\n", encoding="utf-8")
    (INPUT_DIR / "manual_sample.log").write_text(text + "\n", encoding="utf-8")
    (INPUT_DIR / "manual_sample.csv").write_text(
        "type,value\n"
        f"name,{PHRASES[0]}\n"
        f"department,{PHRASES[1]}\n"
        f"code,{PHRASES[2]}\n"
        f"phone,{PHRASES[3]}\n",
        encoding="utf-8",
    )
    (INPUT_DIR / "unsupported_manual.bin").write_text(text + "\n", encoding="utf-8")


def create_docx() -> None:
    document = Document()
    document.add_heading("機密情報マスキングテスト", level=1)
    document.add_paragraph(sample_text())
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "検出語句"
    table.cell(0, 1).text = "値"
    table.cell(1, 0).text = "電話番号"
    table.cell(1, 1).text = PHRASES[3]
    document.save(INPUT_DIR / "manual_sample.docx")


def create_xlsx() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "入力"
    sheet.append(["項目", "値"])
    sheet.append(["担当者", PHRASES[0]])
    sheet.append(["所属", PHRASES[1]])
    sheet.append(["担当コード", PHRASES[2]])
    sheet.append(["電話番号", PHRASES[3]])
    workbook.save(INPUT_DIR / "manual_sample.xlsx")


def create_pptx() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    if title:
        title.text = "機密情報マスキングテスト"
    box = slide.shapes.add_textbox(0, 900000, 8000000, 2500000)
    box.text_frame.text = sample_text()
    presentation.save(INPUT_DIR / "manual_sample.pptx")


def create_pdf() -> None:
    path = INPUT_DIR / "manual_sample.pdf"
    document = fitz.open()
    page = document.new_page()
    for index, line in enumerate(sample_text().splitlines()):
        page.insert_text((72, 72 + index * 18), line, fontname="japan", fontsize=12)
    document.save(path)
    document.close()


def main() -> None:
    create_replacement_table()
    create_text_files()
    create_docx()
    create_xlsx()
    create_pptx()
    create_pdf()


if __name__ == "__main__":
    main()
